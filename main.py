from fastapi import FastAPI, Header, BackgroundTasks, HTTPException, Response
from pydantic import BaseModel
from pptx import Presentation
from sendgrid import SendGridAPIClient
from sendgrid.helpers.mail import Mail, Attachment, FileContent, FileName, FileType, Disposition
import base64, uuid, tempfile, os

import httpx
from fastapi import Request

PUBMED_BASE = "https://eutils.ncbi.nlm.nih.gov/entrez/eutils"

app = FastAPI()
jobs = {}
class PPTReq(BaseModel):
    pmids: list[str] | None = None
    email: str

@app.post("/generatePptAndSend")
async def generate(req: PPTReq,
                   authorization: str = Header(...),
                   bt: BackgroundTasks = BackgroundTasks()):
    job_id = str(uuid.uuid4())
    jobs[job_id] = {"status": "queued"}
    bt.add_task(worker, job_id, req.pmids or [], req.email, authorization)
    return {"jobId": job_id, "message": "started"}

def worker(job_id, pmids, email, auth):
    try:
        print("WORKER START", job_id, email)          # ←1
        jobs[job_id] = {"status": "running"}
        token = os.getenv("SENDGRID_API_KEY") or auth.split()[-1]
        print("TOKEN_BEGINS_WITH", token[:10])        # ←2

        # （PPT 作成処理はそのまま）
        sg = SendGridAPIClient(api_key=token)
        response = sg.send(msg)
        print("SENDGRID STATUS", response.status_code, response.body)  # ←3

        jobs[job_id]["status"] = "finished"
    except Exception as e:
        print("WORKER ERROR", e)                      # ←4
        jobs[job_id] = {"status": "error", "error": str(e)}
    
    jobs[job_id]["status"] = "running"
    prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "自動生成 PPT"
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as f:
        prs.save(f.name); file_path = f.name
    sg = SendGridAPIClient(api_key=auth.split()[-1])
    with open(file_path, "rb") as f:
        data = base64.b64encode(f.read()).decode()
    attachment = Attachment(FileContent(data), FileName("paper.pptx"),
                            FileType("application/vnd.openxmlformats-officedocument.presentationml.presentation"),
                            Disposition("attachment"))
    msg = Mail(from_email=os.getenv("FROM_EMAIL", "no-reply@example.com"),
               to_emails=email, subject="自動生成論文スライド",
               plain_text_content="PPT を添付しました。")
    msg.attachment = attachment
    sg.send(msg)
    jobs[job_id] = {"status": "finished"}

@app.get("/jobs/{job_id}")
def job(job_id: str):
    return jobs.get(job_id) or HTTPException(404, "not found")

@app.api_route("/pubmed/{path:path}", methods=["GET"])
async def proxy_pubmed(path: str, request: Request):
    """
    ChatGPT から来た /pubmed/xxx をそのまま PubMed E-Utilities へ転送して
    レスポンスを返すだけの超シンプルなプロキシ。
    """
    params = dict(request.query_params)
    async with httpx.AsyncClient() as client:
        r = await client.get(f"{PUBMED_BASE}/{path}", params=params, timeout=30)
    # content-type をそのまま引き継ぐ
    return Response(content=r.content, status_code=r.status_code,
                    media_type=r.headers.get("content-type"))

